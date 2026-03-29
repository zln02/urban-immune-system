"""
Urban Immune System — 캡스톤 디자인 발표자료 생성기
원본 Canva 디자인 정밀 재현 (표지 + 감사합니다 2장)
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ── 색상 팔레트 (원본 캡처 기반) ──
DARK_NAVY = RGBColor(0x0B, 0x2A, 0x45)
MID_NAVY = RGBColor(0x14, 0x3D, 0x5C)
LIGHT_TEAL = RGBColor(0x5B, 0xC0, 0xBE)
VERY_LIGHT_TEAL = RGBColor(0xD0, 0xEC, 0xEB)
PALE_GRAY = RGBColor(0xE0, 0xE4, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF4, 0xF6, 0xF9)
TAG_BG = RGBColor(0xE5, 0xF4, 0xF4)
GRAY_TEXT = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY_TEXT = RGBColor(0x9A, 0xA5, 0xB4)
GRID_COLOR = RGBColor(0xE2, 0xE6, 0xEB)
CARD_BORDER = RGBColor(0xE8, 0xEB, 0xEE)
TEAL_DOT = RGBColor(0x4E, 0xC6, 0xC1)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# ── 마진 ──
M_LEFT = Emu(700000)
M_RIGHT = Emu(700000)
CONTENT_W = SLIDE_W - M_LEFT - M_RIGHT


def _set_shape_opacity(shape, alpha_pct):
    """도형에 투명도 설정 (0~100, 100=완전 투명)."""
    spPr = shape._element.spPr
    fill_elem = spPr.find(qn('a:solidFill'))
    if fill_elem is None:
        fill_elem = spPr.find(qn('a:gradFill'))
    if fill_elem is None:
        return
    srgb = fill_elem.find(qn('a:srgbClr'))
    if srgb is None:
        return
    # alpha: 0 = 완전투명, 100000 = 완전불투명
    alpha_val = str(int((100 - alpha_pct) * 1000))
    existing = srgb.find(qn('a:alpha'))
    if existing is not None:
        srgb.remove(existing)
    alpha_elem = srgb.makeelement(qn('a:alpha'), {'val': alpha_val})
    srgb.append(alpha_elem)


def _add_circle(slide, left, top, size, color, opacity_pct=90):
    """반투명 원 장식."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    _set_shape_opacity(shape, opacity_pct)
    return shape


def _add_ring(slide, left, top, size, color, line_w=Pt(2.5), opacity_pct=85):
    """테두리만 있는 원 (링)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = line_w
    # 링 투명도는 line에 적용
    ln = shape._element.spPr.find(qn('a:ln'))
    if ln is not None:
        sf = ln.find(qn('a:solidFill'))
        if sf is not None:
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_val = str(int((100 - opacity_pct) * 1000))
                alpha_elem = srgb.makeelement(qn('a:alpha'), {'val': alpha_val})
                srgb.append(alpha_elem)
    return shape


def _add_grid_lines(slide):
    """배경에 연한 격자 패턴."""
    # 세로선
    for x in range(0, 12192000, 610000):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(x), Emu(0), Emu(8000), SLIDE_H
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.line.fill.background()
        _set_shape_opacity(line, 60)

    # 가로선
    for y in range(0, 6858000, 610000):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(y), SLIDE_W, Emu(8000)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_COLOR
        line.line.fill.background()
        _set_shape_opacity(line, 60)


def _add_decorative_circles(slide, variant="cover"):
    """장식 원·링 배치 (원본과 동일한 위치)."""
    if variant == "cover":
        # 좌하단 큰 원 (반쯤 잘림)
        _add_circle(slide, Emu(-600000), Emu(4600000), Emu(2800000),
                     LIGHT_TEAL, opacity_pct=92)
        _add_ring(slide, Emu(-800000), Emu(4400000), Emu(3200000),
                  LIGHT_TEAL, Pt(2), opacity_pct=88)
        # 우상단
        _add_circle(slide, Emu(10600000), Emu(-900000), Emu(2400000),
                     LIGHT_TEAL, opacity_pct=93)
        _add_ring(slide, Emu(10400000), Emu(-1100000), Emu(2800000),
                  LIGHT_TEAL, Pt(2), opacity_pct=90)
        # 우하단 작은 원
        _add_circle(slide, Emu(10000000), Emu(5400000), Emu(1900000),
                     LIGHT_TEAL, opacity_pct=93)
    else:  # thankyou
        # 좌하단
        _add_circle(slide, Emu(-700000), Emu(4200000), Emu(3000000),
                     LIGHT_TEAL, opacity_pct=92)
        _add_ring(slide, Emu(-900000), Emu(4000000), Emu(3400000),
                  LIGHT_TEAL, Pt(2), opacity_pct=88)
        # 우상단
        _add_circle(slide, Emu(10400000), Emu(-1000000), Emu(2600000),
                     LIGHT_TEAL, opacity_pct=93)
        _add_ring(slide, Emu(10200000), Emu(-1200000), Emu(3000000),
                  LIGHT_TEAL, Pt(2), opacity_pct=90)
        # 우하단
        _add_circle(slide, Emu(10200000), Emu(5200000), Emu(2100000),
                     LIGHT_TEAL, opacity_pct=93)
        _add_ring(slide, Emu(10000000), Emu(5000000), Emu(2500000),
                  LIGHT_TEAL, Pt(2), opacity_pct=90)


def _setup_slide_bg(slide, variant="cover"):
    """배경 + 그리드 + 장식 원."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NEAR_WHITE
    _add_grid_lines(slide)
    _add_decorative_circles(slide, variant)


def _add_text_box(slide, left, top, width, height, text, font_size=12,
                  color=DARK_NAVY, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name="맑은 고딕", italic=False, spacing_after=0,
                  anchor=MSO_ANCHOR.TOP):
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    # anchor (vertical alignment)
    anchor_val = {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't')
    txBox.text_frame._txBody.attrib['anchor'] = anchor_val
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    if spacing_after:
        p.space_after = Pt(spacing_after)
    return txBox


def _add_multiline_text(slide, left, top, width, height, lines,
                        font_name="맑은 고딕", alignment=PP_ALIGN.LEFT):
    """여러 줄 텍스트 (각 줄마다 서식 지정).
    lines: [(text, font_size, color, bold), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox


def _add_tag_pill(slide, left, top, label, width=None):
    """알약 형태 태그 (●  텍스트)."""
    w = width or Emu(len(label) * 160000 + 500000)
    h = Emu(330000)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TAG_BG
    shape.line.fill.background()
    # 라운딩
    shape.adjustments[0] = 0.3

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(100000)
    tf.margin_right = Emu(100000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)

    p = tf.paragraphs[0]
    # 틸 도트 + 텍스트
    run_dot = p.add_run()
    run_dot.text = "●  "
    run_dot.font.size = Pt(8)
    run_dot.font.color.rgb = TEAL_DOT
    run_dot.font.bold = False
    run_dot.font.name = "맑은 고딕"

    run_txt = p.add_run()
    run_txt.text = label
    run_txt.font.size = Pt(9)
    run_txt.font.color.rgb = MID_NAVY
    run_txt.font.bold = True
    run_txt.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    return shape


def _add_line(slide, left, top, width, color=PALE_GRAY, thickness=Emu(12000)):
    """수평 구분선."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, thickness
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_card(slide, left, top, width, height, border_color=CARD_BORDER):
    """흰색 카드 (라운드 사각형)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.05
    return shape


def _add_team_member(slide, left, top, name, role):
    """팀원 이름 + 역할 (세로 배치)."""
    # 이름
    _add_text_box(slide, left, top, Emu(1900000), Emu(280000),
                  name, font_size=13, color=DARK_NAVY, bold=True)
    # 역할 (여러 줄 가능)
    role_lines = role.split('\n')
    y_offset = Emu(300000)
    for rl in role_lines:
        _add_text_box(slide, left, top + y_offset, Emu(1900000), Emu(200000),
                      rl, font_size=8, color=GRAY_TEXT, bold=False)
        y_offset += Emu(160000)


def _add_section_label(slide, left, top, text):
    """섹션 라벨 (PROJECT TEAM, CAPSTONE 등)."""
    _add_text_box(slide, left, top, Emu(3000000), Emu(220000),
                  text, font_size=7, color=LIGHT_GRAY_TEXT, bold=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  슬라이드 1: 표지
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def build_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_slide_bg(slide, "cover")

    # ── 상단 태그 ──
    tag_box = _add_text_box(
        slide, M_LEFT, Emu(380000), Emu(6000000), Emu(280000),
        "", font_size=9, color=LIGHT_TEAL, bold=True,
    )
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    run_dot = p.add_run()
    run_dot.text = "●  "
    run_dot.font.size = Pt(8)
    run_dot.font.color.rgb = TEAL_DOT
    run_dot.font.name = "맑은 고딕"
    run_label = p.add_run()
    run_label.text = "CAPSTONE DESIGN · 동신대학교 컴퓨터공학과"
    run_label.font.size = Pt(9)
    run_label.font.color.rgb = LIGHT_TEAL
    run_label.font.bold = True
    run_label.font.name = "맑은 고딕"

    # ── 메인 제목 ──
    _add_text_box(
        slide, M_LEFT, Emu(900000), Emu(8000000), Emu(700000),
        "도시 면역 체계",
        font_size=44, color=DARK_NAVY, bold=True,
    )

    # ── 영문 부제 ──
    _add_text_box(
        slide, M_LEFT, Emu(1650000), Emu(8000000), Emu(400000),
        "URBAN IMMUNE SYSTEM",
        font_size=16, color=MID_NAVY, bold=False,
    )

    # ── 설명 카드 ──
    card = _add_card(slide, M_LEFT, Emu(2250000), Emu(4800000), Emu(400000))
    # 카드 위에 텍스트
    desc_box = _add_text_box(
        slide, Emu(M_LEFT + Emu(180000)), Emu(2250000),
        Emu(4500000), Emu(400000),
        "AI 기반 감염병 조기경보 시스템",
        font_size=12, color=MID_NAVY, bold=False,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ── 3-Layer 태그 ──
    tags = ["약국 OTC", "하수 바이오마커", "검색 트렌드"]
    tag_x = M_LEFT
    for t in tags:
        _add_tag_pill(slide, tag_x, Emu(2850000), t)
        tag_x += Emu(2300000)

    # ── 구분선 1 ──
    _add_line(slide, M_LEFT, Emu(3500000), CONTENT_W)

    # ── PROJECT TEAM 라벨 ──
    _add_section_label(slide, M_LEFT, Emu(3620000), "PROJECT TEAM")

    # ── 팀원 5명 ──
    members = [
        ("박진영", "PM / ML Lead\n아이디어 설계 · 3-Layer 교차검증"),
        ("이경준", "Backend\n서버 · API 개발"),
        ("이우형", "Data Engineer\n데이터 파이프라인"),
        ("김나영", "Frontend\n대시보드 UI/UX"),
        ("박정빈", "DevOps / QA\n배포 · 테스트"),
    ]
    x_start = M_LEFT
    x_gap = Emu(2150000)
    for i, (name, role) in enumerate(members):
        _add_team_member(slide, x_start + x_gap * i, Emu(3880000), name, role)

    # ── 구분선 2 ──
    _add_line(slide, M_LEFT, Emu(4700000), CONTENT_W)

    # ── CAPSTONE 섹션 (좌측) ──
    _add_section_label(slide, M_LEFT, Emu(4820000), "CAPSTONE")

    _add_text_box(
        slide, M_LEFT, Emu(5050000), Emu(5000000), Emu(280000),
        "동신대학교 컴퓨터공학과",
        font_size=11, color=DARK_NAVY, bold=True,
    )
    _add_text_box(
        slide, M_LEFT, Emu(5310000), Emu(5000000), Emu(230000),
        "캡스톤 디자인 · 2026",
        font_size=9, color=GRAY_TEXT, bold=False,
    )

    # ── 수상 뱃지 (우측) ──
    badge_left = Emu(7200000)
    badge_w = Emu(4300000)
    badge_h = Emu(550000)
    badge_card = _add_card(slide, badge_left, Emu(4900000), badge_w, badge_h)

    badge_txt = slide.shapes.add_textbox(
        badge_left, Emu(4900000), badge_w, badge_h
    )
    tf = badge_txt.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000)
    tf.margin_right = Emu(150000)
    # 세로 중앙
    badge_txt.text_frame._txBody.attrib['anchor'] = 'ctr'

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run_icon = p.add_run()
    run_icon.text = "★  "
    run_icon.font.size = Pt(11)
    run_icon.font.color.rgb = RGBColor(0xD4, 0xA5, 0x37)
    run_icon.font.bold = True
    run_icon.font.name = "맑은 고딕"

    run_main = p.add_run()
    run_main.text = "제1회 AI 아이디어 공모전 대상 수상작 기반"
    run_main.font.size = Pt(10)
    run_main.font.color.rgb = MID_NAVY
    run_main.font.bold = True
    run_main.font.name = "맑은 고딕"

    # ── 하단 저작권 라인 ──
    _add_text_box(
        slide, Emu(0), Emu(6500000), SLIDE_W, Emu(200000),
        "© 2026 Urban Immune System Team · 동신대학교",
        font_size=7, color=LIGHT_GRAY_TEXT, bold=False,
        alignment=PP_ALIGN.CENTER,
    )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  슬라이드 2: 감사합니다
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def build_slide_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _setup_slide_bg(slide, "thankyou")

    # ── 상단 태그 ──
    tag_box = _add_text_box(
        slide, Emu(0), Emu(700000), SLIDE_W, Emu(280000),
        "", font_size=9, color=LIGHT_TEAL, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_dot = p.add_run()
    run_dot.text = "●  "
    run_dot.font.size = Pt(8)
    run_dot.font.color.rgb = TEAL_DOT
    run_dot.font.name = "맑은 고딕"
    run_label = p.add_run()
    run_label.text = "URBAN IMMUNE SYSTEM"
    run_label.font.size = Pt(9)
    run_label.font.color.rgb = LIGHT_TEAL
    run_label.font.bold = True
    run_label.font.name = "맑은 고딕"

    # ── 감사합니다 ──
    _add_text_box(
        slide, Emu(0), Emu(1600000), SLIDE_W, Emu(1000000),
        "감사합니다",
        font_size=54, color=DARK_NAVY, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # ── Thank You ──
    _add_text_box(
        slide, Emu(0), Emu(2700000), SLIDE_W, Emu(450000),
        "Thank You",
        font_size=18, color=MID_NAVY, bold=False,
        alignment=PP_ALIGN.CENTER, italic=True,
    )

    # ── 구분선 (중앙) ──
    line_w = Emu(5500000)
    line_left = Emu((12192000 - 5500000) // 2)
    _add_line(slide, line_left, Emu(3500000), line_w)

    # ── 팀원 5명 (중앙 정렬) ──
    members = [
        ("박진영", "PM / ML Lead"),
        ("이경준", "Backend"),
        ("이우형", "Data Eng"),
        ("김나영", "Frontend"),
        ("박정빈", "DevOps / QA"),
    ]
    total_members = len(members)
    member_w = Emu(1800000)
    total_w = member_w * total_members + Emu(200000) * (total_members - 1)
    x_start = Emu((12192000 - total_w) // 2)
    x_gap = Emu(2000000)

    for i, (name, role) in enumerate(members):
        x = x_start + x_gap * i
        # 이름
        _add_text_box(
            slide, x, Emu(3700000), member_w, Emu(280000),
            name, font_size=13, color=DARK_NAVY, bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        # 역할
        _add_text_box(
            slide, x, Emu(3980000), member_w, Emu(200000),
            role, font_size=8, color=GRAY_TEXT, bold=False,
            alignment=PP_ALIGN.LEFT,
        )

    # ── 프로젝트 정보 카드 ──
    info_card_w = Emu(7200000)
    info_card_h = Emu(1300000)
    info_card_left = Emu((12192000 - 7200000) // 2)
    info_card_top = Emu(4500000)
    _add_card(slide, info_card_left, info_card_top, info_card_w, info_card_h)

    # 카드 제목
    _add_text_box(
        slide, info_card_left + Emu(250000), info_card_top + Emu(80000),
        Emu(6700000), Emu(250000),
        "Urban Immune System — AI 기반 감염병 조기경보 시스템",
        font_size=9, color=DARK_NAVY, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    _add_text_box(
        slide, info_card_left + Emu(250000), info_card_top + Emu(310000),
        Emu(6700000), Emu(200000),
        "동신대학교 컴퓨터공학과 캡스톤 디자인",
        font_size=8, color=GRAY_TEXT, bold=False,
        alignment=PP_ALIGN.LEFT,
    )

    # 링크 정보
    links = [
        ("📦 GitHub", "github.com/xxx/urban-immune-system"),
        ("📋 Notion", "notion.so/xxx (프로젝트 문서)"),
        ("🚀 Demo", "34.64.122.238:8501 (프로토타입)"),
    ]
    link_y = info_card_top + Emu(530000)
    for icon_label, url in links:
        link_box = slide.shapes.add_textbox(
            info_card_left + Emu(250000), link_y, Emu(6700000), Emu(200000)
        )
        ltf = link_box.text_frame
        ltf.word_wrap = False
        lp = ltf.paragraphs[0]
        r1 = lp.add_run()
        r1.text = f"{icon_label}   "
        r1.font.size = Pt(8)
        r1.font.color.rgb = MID_NAVY
        r1.font.bold = True
        r1.font.name = "맑은 고딕"
        r2 = lp.add_run()
        r2.text = url
        r2.font.size = Pt(8)
        r2.font.color.rgb = GRAY_TEXT
        r2.font.bold = False
        r2.font.name = "맑은 고딕"
        link_y += Emu(210000)

    # ── 팀원 요약 (카드 아래) ──
    _add_text_box(
        slide, Emu(0), Emu(5950000), SLIDE_W, Emu(220000),
        "박진영(PM/ML) · 이경준(Backend) · 이우형(Data) · 김나영(Frontend) · 박정빈(DevOps)",
        font_size=8, color=GRAY_TEXT, bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # ── 하단 정보 ──
    _add_text_box(
        slide, Emu(0), Emu(6200000), SLIDE_W, Emu(220000),
        "동신대학교 컴퓨터공학과 · 캡스톤 디자인 2026",
        font_size=9, color=GRAY_TEXT, bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # ── 하단 저작권 ──
    _add_text_box(
        slide, Emu(0), Emu(6450000), SLIDE_W, Emu(200000),
        "© 2026 Urban Immune System Team · 동신대학교",
        font_size=7, color=LIGHT_GRAY_TEXT, bold=False,
        alignment=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_cover(prs)
    build_slide_thankyou(prs)

    out_path = os.path.join(
        os.path.dirname(__file__),
        "Urban_Immune_System_캡스톤디자인.pptx",
    )
    prs.save(out_path)
    print(f"생성 완료: {out_path}")


if __name__ == "__main__":
    main()
