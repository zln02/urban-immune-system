"""HTML deck → PPTX export.

Playwright(chromium)로 각 슬라이드 1920×1080 PNG 캡처 후
python-pptx 16:9 빈 슬라이드에 풀스크린 박는다.

산출물: docs/slides/2026-04-30_중간발표.pptx (24장)
"""
from __future__ import annotations

import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

PROJECT = Path(__file__).parent.parent
DECK_HTML = PROJECT / "frontend" / "public" / "slides" / "index.html"
OUT_PPTX = PROJECT / "docs" / "slides" / "2026-04-30_중간발표.pptx"
TMP_DIR = PROJECT / "docs" / "slides" / "_export_png"
TMP_DIR.mkdir(parents=True, exist_ok=True)

W, H = 1920, 1080
N_SLIDES = 24


def capture_pngs() -> list[Path]:
    out_paths: list[Path] = []
    deck_url = f"file://{DECK_HTML.resolve()}"
    with sync_playwright() as p:
        browser = p.chromium.launch()
        context = browser.new_context(viewport={"width": W, "height": H}, device_scale_factor=1)
        page = context.new_page()
        for i in range(1, N_SLIDES + 1):
            url = f"{deck_url}#{i}"
            page.goto(url, wait_until="networkidle")
            # 폰트·SVG 렌더 안정화 대기
            page.wait_for_timeout(800)
            png = TMP_DIR / f"slide_{i:02d}.png"
            page.screenshot(path=str(png), full_page=False, omit_background=False)
            out_paths.append(png)
            print(f"  ✓ {png.name}")
        browser.close()
    return out_paths


def build_pptx(pngs: list[Path]) -> None:
    prs = Presentation()
    # 표준 widescreen 16:9 — 13.333 × 7.5 inch (PowerPoint default)
    prs.slide_width = Emu(12192000)   # 13.333 inch
    prs.slide_height = Emu(6858000)   # 7.5 inch
    blank_layout = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height,
        )
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))


def main() -> int:
    if not DECK_HTML.exists():
        print(f"❌ deck not found: {DECK_HTML}")
        return 1
    print(f"▶ Capturing {N_SLIDES} slides {W}×{H}")
    pngs = capture_pngs()
    print(f"▶ Building PPTX → {OUT_PPTX}")
    build_pptx(pngs)
    size_kb = OUT_PPTX.stat().st_size // 1024
    print(f"✅ done · {OUT_PPTX} ({size_kb} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
