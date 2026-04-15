"""Urban Immune System — 팀 킥오프 PPTX 자동 생성.

실행:
    python scripts/gen_kickoff_pptx.py

출력:
    docs/meeting-notes/2026-04-15_팀킥오프.pptx

변환(선택):
    libreoffice --headless --convert-to pdf \\
      docs/meeting-notes/2026-04-15_팀킥오프.pptx \\
      --outdir docs/meeting-notes/
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "meeting-notes" / "2026-04-15_팀킥오프.pptx"

NAVY = RGBColor(0x1F, 0x2D, 0x5B)
ACCENT = RGBColor(0xE6, 0x3B, 0x3B)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)


def add_title_slide(prs, title: str, subtitle: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 네이비 배경 띠
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(1),
             title, pt=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(0.6),
             subtitle, pt=20, color=GRAY)


def add_content_slide(prs, title: str, lines: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             title, pt=32, bold=True, color=NAVY)
    # 구분선
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15),
                              Inches(12), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    body = "\n".join(f"• {l}" for l in lines)
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5),
             body, pt=18, color=GRAY)


def add_table_slide(prs, title: str, header: list[str], rows: list[list[str]]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             title, pt=32, bold=True, color=NAVY)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15),
                              Inches(12), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    cols = len(header)
    rows_count = len(rows) + 1
    tbl_shape = s.shapes.add_table(rows_count, cols, Inches(0.6), Inches(1.4),
                                   Inches(12), Inches(5.0))
    table = tbl_shape.table
    for c, h in enumerate(header):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for ri, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = table.cell(ri, c)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = GRAY
            if ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT


def add_text(slide, x, y, w, h, text, pt=18, bold=False, color=GRAY) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, para in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        for r in p.runs:
            r.font.size = Pt(pt)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "NanumGothic"


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    add_title_slide(
        prs,
        "Urban Immune System — 팀 킥오프",
        "2026-04-15 · 감염병 조기경보 AI · LG전자 DX School 공모전 대상(1등)",
    )

    # 2. 프로젝트 한 줄 요약
    add_content_slide(prs, "🎯 프로젝트 한 줄 요약", [
        "3-Layer 비의료 신호(약국 OTC + 하수 바이오마커 + 검색 트렌드)로",
        "감염병 1~3주 선행 감지 + RAG-LLM 자동 경보 리포트",
        "B2G 납품 지향 — 질병관리청·광역지자체 레퍼런스 확보 목표",
        "기술: TFT + Autoencoder + FastAPI + TimescaleDB + Streamlit/Next.js",
    ])

    # 3. 현재 완성도 (정직한 평가)
    add_table_slide(prs, "📊 현재 완성도 (L1 스켈레톤 ~ L5 프로덕션)",
        ["모듈", "레벨", "요약"],
        [
            ["src/ Streamlit", "L3", "5탭 렌더, 시뮬 데이터, 일부 하드코딩"],
            ["backend/ FastAPI", "L2", "placeholder 다수, JWT·감사로그 미구현"],
            ["pipeline/", "L2", "Producer 4개, Kafka Consumer 없음"],
            ["ml/ TFT/AE/RAG", "L2", "모델 정의만, 체크포인트·서빙 wiring 미완"],
            ["frontend/ Next.js", "L1", "컴포넌트 스텁"],
            ["tests/", "trivial", "25 pass, 커버리지 낮음"],
        ])

    # 4. 팀 역할
    add_table_slide(prs, "👥 팀 역할 + 담당 모듈",
        ["팀원", "역할", "담당 CLAUDE.md"],
        [
            ["박진영", "PM / ML Lead", "ml/, docs/, docs/business/"],
            ["이경준", "Backend", "backend/"],
            ["이우형", "Data Engineer", "pipeline/"],
            ["김나영", "Frontend", "src/ (Phase1), frontend/ (Phase2)"],
            ["박정빈", "DevOps / QA", "infra/, .github/, tests/"],
        ])

    # 5. GitHub 협업 규칙
    add_content_slide(prs, "🌿 GitHub 협업 규칙", [
        "브랜치: main ← develop ← feature/<이니셜>-<작업명>",
        "PR: develop 1인 승인 / main 2인 승인 + 6 CI Job 전부 통과",
        "CODEOWNERS — 모듈별 자동 리뷰어 (.github/CODEOWNERS)",
        "ISSUE/PR 템플릿 (.github/ISSUE_TEMPLATE, pull_request_template)",
        "Commit 접두어: feat / fix / docs / chore / decision / troubleshoot / milestone",
    ])

    # 6. 권한 표
    add_table_slide(prs, "🔐 팀원 GitHub 권한",
        ["팀원", "GitHub Role", "주요 권한"],
        [
            ["박진영", "Admin", "모든 것 (브랜치 보호·Secrets·릴리스)"],
            ["이경준", "Write", "자기 feature push / PR 생성·리뷰"],
            ["이우형", "Write", "동일"],
            ["김나영", "Write", "동일"],
            ["박정빈", "Maintain", "Actions 재실행·환경변수·PR merge"],
        ])

    # 7. 환경 세팅 3분
    add_content_slide(prs, "⚡ 환경 세팅 3분 가이드", [
        "ssh wlsdud5035@34.64.124.90",
        "cd ~/urban-immune-system && git fetch origin",
        "git switch -c feature/<이니셜>-<작업명> origin/develop",
        "cd <모듈> && claude     ← 배지가 뜨면 성공",
        "상세: docs/meeting-notes/setup-per-role.md",
    ])

    # 8. CLI 일상 사용법
    add_content_slide(prs, "🤖 Claude Code 일상 사용", [
        "cd <모듈> && claude — 해당 CLAUDE.md 자동 로드 + 역할 배지",
        "/commit /review-pr /simplify 내장 Skills",
        "Stop hook: 세션 종료 시 memory + portfolio 자동 기록",
        "Opus 지휘 → Haiku/Sonnet 서브에이전트에 병렬 위임(3개 한도)",
    ])

    # 9. 주차별 마일스톤
    add_table_slide(prs, "📅 캡스톤 발표까지 마일스톤",
        ["주차", "목표", "담당"],
        [
            ["7 (다음주)", "P0 수치 정직성 복구 (F1/MCC/AUPRC 실측)", "박진영"],
            ["8", "P1 멀티 시즌 검증 (2022-23, 2023-24)", "박진영"],
            ["8-9", "P2 실데이터 E2E (Prefect · DB 적재)", "이우형·이경준"],
            ["9-10", "P3 ML 추론 wiring (ml/serve.py · LGBM)", "박진영"],
            ["10-11", "P4 라이브 대시보드 · 주간 리포트", "김나영"],
            ["11", "B2G 산출물 초안 · 발표 리허설", "전원"],
            ["12 (발표)", "데모 스크립트 + 1-click 기동 검증", "전원"],
        ])

    # 10. 1등 진입 3대 핵심
    add_content_slide(prs, "🏆 1등 진입 3대 핵심", [
        "(1) 수치 정직성 복구 — README F1/Precision 값이 실제 계산 코드로 재현 가능해야 함",
        "(2) 멀티 시즌 검증 — 2022-23, 2023-24 추가 → 과적합 방어",
        "(3) 라이브 데모 MVP — docker compose 1명령 15분 내 기동",
        "→ 대학 캡스톤 1등 확률 65% → 85% 목표",
    ])

    # 11. 기술 스택 변경 제안
    add_table_slide(prs, "🔧 기술 스택 변경 제안",
        ["영역", "현재", "제안"],
        [
            ["파이프라인", "Kafka KRaft", "Prefect Cloud 무료 (오버엔지니어링 제거)"],
            ["배포", "K8s GKE", "Phase1은 Cloud Run / Next.js는 Vercel"],
            ["ML", "TFT 단독", "LightGBM 베이스라인 → TFT 단계 승급"],
            ["벡터 DB", "Qdrant", "pgvector (TimescaleDB 확장)"],
            ["RAG", "LangChain", "직접 구현 or LlamaIndex"],
        ])

    # 12. 병행 도구
    add_content_slide(prs, "🛠 병행 도구 (Claude Code 외)", [
        "GitHub Copilot — 학생 무료, 전원 즉시 활성화 권장",
        "Cursor 무료 티어 (2000 completions/월) — 복잡 멀티파일",
        "Notion — 회의록·API 명세·데이터 계약",
        "Discord + GitHub webhook — 실시간 알림 (Slack 90일 제한 회피)",
        "Figma — Phase2 UI 설계 · Otter.ai — 회의 자동 전사",
    ])

    # 13. 포트폴리오 자동 메모리
    add_content_slide(prs, "🧠 포트폴리오 자동 메모리 (신규)", [
        "docs/portfolio/ — timeline / decisions / troubleshooting / milestones / retrospectives",
        "Stop hook 이 세션 종료 시 timeline.md 에 자동 append",
        "커밋 접두어로 자동 분류: decision:/troubleshoot:/milestone:",
        "캡스톤 발표 때 portfolio.html 시연 → 시행착오·결정 과정 심사자에게 제시",
    ])

    # 14. 오늘 액션 아이템
    add_content_slide(prs, "✅ 오늘 액션 아이템", [
        "박진영: Collaborator 4명 초대 + Secrets 7개 등록 + Branch protection",
        "전원: 초대 수락 → SSH 접속 → cd <모듈> && claude 테스트",
        "박정빈: Frontend Lint CI 수정 (package-lock.json 이미 생성됨)",
        "박진영: P0 성능 측정 노트북 착수 (D+7 데드라인)",
        "전원: GitHub Copilot 학생 인증 (education.github.com)",
        "전원: decision: 또는 troubleshoot: 접두어 커밋 1건씩 체험",
    ])

    # 15. Q&A
    add_content_slide(prs, "💬 Q&A", [
        "자유 질문 / 역할 조정 / 기술 스택 합의",
        "다음 정기 회의 시간 결정",
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"✅ {OUT.relative_to(OUT.parents[2])} 생성 ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} 슬라이드)")


if __name__ == "__main__":
    build()
