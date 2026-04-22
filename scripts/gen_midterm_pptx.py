"""중간발표 PPTX 자동 생성 (2026-04-30).

실행: python scripts/gen_midterm_pptx.py
출력: docs/slides/2026-04-30_중간발표.pptx

구조: (1) 어느정도 만들어졌나 → (2) 처음 방향에서 추가/제외 → (3) 최종 제품 방향
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "slides" / "2026-04-30_중간발표.pptx"

NAVY = RGBColor(0x1F, 0x2D, 0x5B)
ACCENT = RGBColor(0xE6, 0x3B, 0x3B)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0x9E, 0x73)
ORANGE = RGBColor(0xD5, 0x5E, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)


def text(s, x, y, w, h, body, pt=16, bold=False, color=GRAY, name="NanumGothic"):
    box = s.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = ln
        for r in p.runs:
            r.font.size = Pt(pt); r.font.bold = bold; r.font.color.rgb = color; r.font.name = name


def title_slide(prs, t, sub, badge=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(2.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(1), t, pt=40, bold=True, color=WHITE)
    text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.8), sub, pt=20, color=RGBColor(0xCC, 0xD5, 0xE0))
    if badge:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.8),
                                  Inches(12), Inches(0.9))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7); card.line.color.rgb = ORANGE
        text(s, Inches(0.9), Inches(2.95), Inches(11.4), Inches(0.7), badge, pt=18, bold=True, color=NAVY)


def header(s, t):
    text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), t, pt=28, bold=True, color=NAVY)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.1), Inches(12), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()


def content(prs, t, lines, pt=15):
    s = prs.slides.add_slide(prs.slide_layouts[6]); header(s, t)
    text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(5.8),
         "\n".join(f"• {l}" for l in lines), pt=pt, color=GRAY)


def table_slide(prs, t, h, rows, fs=12, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); header(s, t)
    cols = len(h); rc = len(rows) + 1
    ht = Inches(5.0) if note else Inches(5.5)
    tbl = s.shapes.add_table(rc, cols, Inches(0.6), Inches(1.35), Inches(12), ht).table
    for c, v in enumerate(h):
        cell = tbl.cell(0, c); cell.text = v
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(fs + 1); r.font.color.rgb = WHITE; r.font.name = "NanumGothic"
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for ri, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = tbl.cell(ri, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = GRAY; r.font.name = "NanumGothic"
            if ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
    if note:
        text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.7), note, pt=12, color=RGBColor(0x66, 0x66, 0x66))


def two_table_slide(prs, t, h1, rows1, h2, rows2, label1="🔴 제외", label2="🟢 추가", fs=11):
    """좌우 2분할 테이블."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); header(s, t)

    text(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4), label1, pt=16, bold=True, color=RED)
    cols1 = len(h1); rc1 = len(rows1) + 1
    tbl1 = s.shapes.add_table(rc1, cols1, Inches(0.5), Inches(1.8), Inches(6.2), Inches(5.3)).table
    for c, v in enumerate(h1):
        cell = tbl1.cell(0, c); cell.text = v
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(fs + 1); r.font.color.rgb = WHITE; r.font.name = "NanumGothic"
        cell.fill.solid(); cell.fill.fore_color.rgb = RED
    for ri, row in enumerate(rows1, start=1):
        for c, v in enumerate(row):
            cell = tbl1.cell(ri, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = GRAY; r.font.name = "NanumGothic"

    text(s, Inches(6.9), Inches(1.3), Inches(6), Inches(0.4), label2, pt=16, bold=True, color=GREEN)
    cols2 = len(h2); rc2 = len(rows2) + 1
    tbl2 = s.shapes.add_table(rc2, cols2, Inches(6.9), Inches(1.8), Inches(6.2), Inches(5.3)).table
    for c, v in enumerate(h2):
        cell = tbl2.cell(0, c); cell.text = v
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(fs + 1); r.font.color.rgb = WHITE; r.font.name = "NanumGothic"
        cell.fill.solid(); cell.fill.fore_color.rgb = GREEN
    for ri, row in enumerate(rows2, start=1):
        for c, v in enumerate(row):
            cell = tbl2.cell(ri, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = GRAY; r.font.name = "NanumGothic"


prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

# Slide 1 — 표지
title_slide(prs, "Urban Immune System",
            "3-Layer 비의료 신호 기반 AI 감염병 조기경보 · 캡스톤 중간발표 · 2026-04-30",
            "🏆 제1회 2026 데이터로 미래를 그리는 AI 아이디어 공모전 대상 (LG전자 DX School)")

# Slide 2 — 문제 정의
content(prs, "🔴 문제 정의 — 왜 어려운가", [
    "기존: 환자가 병원 가야 데이터 생김 → 이미 확산된 후",
    "질병관리청 주간 보고 — 1~2주 지연",
    "의사 신고 — 확진 후 신고",
    "학교 결석 — 증상 발현 후",
    "가설: 병원 가기 '전' 시민의 행동(약국·검색·하수) 에 신호가 있다",
])

# Slide 3 — 3-Layer 아이디어
table_slide(prs, "💡 3-Layer 아이디어 + GFT 실패 회피",
    ["Layer", "데이터 소스", "선행 시간"],
    [
        ["💊 약국 OTC", "네이버 쇼핑인사이트", "~2주"],
        ["🚰 하수 바이오마커", "환경부 KOWAS", "~3주"],
        ["🔍 검색 트렌드", "네이버 데이터랩", "~1주"],
    ],
    note="3개 동시 상승 = 경보. 하수(L2)는 사람 행동과 독립 → Google Flu Trends 실패(단일 신호 2배 과대추정) 회피")

# Slide 4 — 아키텍처
content(prs, "🏗 아키텍처", [
    "[3-Layer 수집] → [Kafka] → [TimescaleDB] → [Autoencoder] → [TFT] → [RAG-LLM] → [대시보드]",
    "Dashboard: Streamlit (Phase1) + Next.js+Deck.gl (Phase2)",
    "Backend: FastAPI async · Pipeline: Kafka KRaft · ML: PyTorch Forecasting",
    "LLM: GPT-4o + Claude Sonnet 4.6 + Qdrant RAG",
    "Infra: Docker + K8s(GKE) + GitHub Actions 11 Job 🟢",
])

# Slide 5 — 🆕 처음 계획 (공모전 제안 시점)
table_slide(prs, "📝 처음 계획 (공모전 2026-03 제안 시점)",
    ["영역", "원래 6개월 로드맵 목표"],
    [
        ["데이터", "심평원 처방약 + K-WBE API + 네이버 검색 실시간"],
        ["ML", "TFT + Autoencoder + Isolation Forest 앙상블"],
        ["UI", "Next.js + Deck.gl 3D 위험도 맵 (프로덕션)"],
        ["배포", "K8s(GKE asia-northeast3) 24/7 운영"],
        ["성능 주장", "F1 0.71 · Precision 1.00 · 오경보 0건"],
        ["데이터 주장", "실데이터 26주 분석 완료"],
        ["RAG", "Qdrant 벡터DB 실 임베딩"],
        ["법 대응", "언급 없음"],
    ],
    note="솔직한 인정: 공모전 수치는 단일 임계값 최적화 후 가장 좋은 값만 표기. Precision=1.00은 클래스 불균형에서 자동 발생. 실데이터는 심평원/K-WBE API 미접근 — 합성 데이터였음.", fs=12)

# Slide 6 — 🆕 초기 → 현재: 추가/제외 (좌우 분할)
two_table_slide(prs, "🔄 초기 → 현재: 추가/제외된 것",
    ["기능", "대안"],
    [
        ["K8s GKE 배포", "Docker Compose · GKE 설계만"],
        ["Next.js 프로덕션 UI", "Phase 2 연기"],
        ["Qdrant RAG 실 임베딩", "하드코딩 RAG 데모"],
        ["심평원·K-WBE 실 API", "KDCA ILINet CSV"],
        ["Kafka Consumer", "asyncpg 직접 INSERT 대안"],
        ["TFT 실학습 1차", "LightGBM 베이스라인 우선"],
        ["F1 0.71 · Precision 1.00", "walk-forward → F1 0.643 실측"],
    ],
    ["기능", "추가 이유"],
    [
        ["Okabe-Ito 색맹 팔레트", "B2G WCAG AA 필수"],
        ["9개 모듈 CLAUDE.md 에이전트", "팀원 역할 특화"],
        ["포트폴리오 자동 메모리", "면접·B2G 심사"],
        ["실측 재현 노트북", "P0 정직성 복구"],
        ["CI 11 Job (Trivy·CodeQL·Bandit)", "ISMS-P 25%"],
        ["팀원 이슈 #3-6 + CODEOWNERS", "원맨쇼 리스크 해소"],
        ["docs/business/ 5종", "상용화 준비"],
    ],
    fs=10)

# Slide 7 — 현재 완성도
table_slide(prs, "📊 현재 완성도 (솔직 공개)",
    ["모듈", "완성도", "데모", "다음 주 목표"],
    [
        ["Streamlit (src/)", "75%", "✅ 5탭", "스크린샷 5장"],
        ["Backend (backend/)", "45%", "✅ /health", "실 DB 쿼리 #3"],
        ["Pipeline (pipeline/)", "50%", "✅ Producer", "Consumer #4"],
        ["ML (ml/)", "40%", "✅ JSON", "LightGBM 베이스라인"],
        ["Next.js (frontend/)", "30%", "🟡 scaffold", "Phase 2 연기"],
        ["Infra (infra/)", "75%", "✅ CI 11", "Branch Protect #6"],
    ],
    note="핵심 성과 5개: (1)P0 정직성 복구 (2)Okabe-Ito 팔레트 (3)9개 에이전트 (4)팀원 이슈 #3-6 (5)Docker Compose 5서비스 안정", fs=12)

# Slide 8 — 핵심 기술 TFT + RAG
content(prs, "🤖 핵심 기술 — TFT + RAG", [
    "Temporal Fusion Transformer (Google, 2021)",
    "  → 다변량 시계열 + Attention (어느 Layer 가 예측에 기여?)",
    "왜 TFT? (TimesFM·Moirai·Chronos 2024 등장에도)",
    "  → Foundation model 은 공중보건 도메인 검증 미흡",
    "  → TFT 는 의료 도메인 검증 논문 다수 (말라리아·뎅기·독감)",
    "  → B2G 납품 = '왜 경보?' 답변 가능성이 정확도만큼 중요",
    "RAG: 경보 → Qdrant 유사사례 검색 → GPT-4o/Claude 리포트",
    "현재: 하드코딩 RAG 5건 + OpenAI SDK. 실 임베딩은 Phase 3.",
])

# Slide 9 — 실측 성능
table_slide(prs, "📈 재현 가능한 실측 성능 (P0 복구 완료)",
    ["모델", "F1", "MCC", "AUPRC", "FP"],
    [
        ["L1 약국 단독", "0.621", "0.399", "0.820", "11"],
        ["L2 하수 단독", "0.615", "0.359", "0.838", "9"],
        ["L3 검색 단독", "0.737", "0.588", "0.851", "3"],
        ["3-Layer 앙상블", "0.643", "0.442", "0.885", "10"],
    ],
    note="2024-25 합성 데이터(seed=42, 52주). Granger: L1 p=0.0000 · L2 p=0.0000 · L3 p=0.0193 → 3 Layer 모두 p<0.05 유의. 재현: python analysis/notebooks/performance_measurement.py")

# Slide 10 — 솔직 자기 진단
table_slide(prs, "✋ 솔직한 자기 진단",
    ["구분", "공모전 시점", "현재 (정직)"],
    [
        ["F1", "0.71 (하드코딩)", "0.643 실측"],
        ["Precision", "1.00 단독 강조", "MCC 0.442 + AUPRC 0.885 병기"],
        ["Granger", "p<0.05 HTML 텍스트", "실측 p-value 표시"],
        ["데이터", "실데이터 26주 (주장)", "재현 가능한 합성 · KDCA 실 수집 진행 중"],
    ],
    note="'공모전 수치 지키기'보다 '공공기관 납품 시 재현 가능한 숫자' 가 B2G 신뢰의 본질. 약점 인정이 발표 신뢰도 상승.")

# Slide 11 — 경쟁 맵
table_slide(prs, "🌍 경쟁 맵 (2026 Nature Comms·Lancet DH 반영)",
    ["기준", "BlueDot", "CDC NWSS", "KAIST OTC", "Xu 2025", "우리"],
    [
        ["데이터", "뉴스+항공", "하수", "OTC", "하수+SNS+이동", "OTC+하수+검색"],
        ["모델", "ML ens.", "통계", "DL", "TFT", "TFT+AE+RAG"],
        ["설명성", "낮음", "중", "중", "중", "높음 (RAG)"],
        ["한국 특화", "X", "X", "O", "X", "O"],
        ["비용", "고가", "무료", "연구", "연구", "오픈소스"],
    ],
    note="차별점: 약국 OTC 추가 + 한국어 RAG + Attention 해석성. 약해진 주장: 3-source 최초(Xu 존재), TFT SOTA(foundation model 등장) → 재포지셔닝 완료", fs=11)

# Slide 12 — 법·규제
table_slide(prs, "⚖️ 법·규제 대응",
    ["법·규정", "진행률", "이슈"],
    [
        ["ISMS-P", "25%", "감사로그·JWT·의존성스캔 완료"],
        ["개인정보보호법 29조", "55%", "인증 미들웨어 W18"],
        ["네이버 API 재판매 금지", "⚠️", "네이버 공식 문의 W18"],
        ["감염병예방법 11·15조", "⚠️", "면책 고지 + CDSS 로드맵"],
        ["의료기기법 2조", "✅", "비의료기기 (공중보건 감시)"],
    ],
    note="자동 점검: /legal-review 스킬 주 1회 실행")

# Slide 13 — 🆕 최종 제품 비전
table_slide(prs, "🎯 최종 제품 비전 (Phase 5 · 2027 B2G 납품)",
    ["영역", "Phase 1 (지금)", "Phase 5 (납품)"],
    [
        ["데이터", "합성 + 소량 실측", "3-Layer 실시간 자동"],
        ["예측", "LightGBM 베이스라인", "TFT + Autoencoder 앙상블"],
        ["UI", "Streamlit 로컬", "Next.js + PWA 모바일"],
        ["배포", "Docker Compose", "GKE + Cloudflare + 24/7"],
        ["보안", "기본 SAST", "ISMS-P 인증 완료"],
        ["고객", "캡스톤 심사", "질병관리청 · 17개 광역"],
        ["가격", "무료", "지자체 연 2~6천만 · 광역 연 1~3억"],
    ],
    note="타깃 PoC 순위: 🥇질병관리청 감염병감시지원단 🥈서울시 감염병관리지원단 🥉WHO 협력센터")

# Slide 14 — 상용화 로드맵
content(prs, "🗺 상용화 로드맵 (5 Phase)", [
    "Phase 1 (Now 4/30) — 캡스톤 중간발표 MVP (Streamlit)",
    "Phase 2 (5-6월 기말) — 실데이터 E2E · LightGBM 실측 · LLM 호출",
    "Phase 3 (여름) — GKE 배포 · 질병관리청 PoC 제안 · 법인 설립 검토",
    "Phase 4 (가을) — KOSA 신고 · 조달청 혁신제품 지정 신청",
    "Phase 5 (2027+) — ISMS-P 인증 · 첫 유료 계약 · 해외 확장",
    "예상 6개월 결과: 정부과제 2~5천만 + PoC MOU 1건 + 논문 1편",
    "(KDD Health Day 또는 AAAI AI in Healthcare 투고 예정)",
])

# Slide 15 — Live Demo
content(prs, "🎥 Live Demo (5분)", [
    "URL: http://34.64.124.90:8501 (GCP 방화벽 허용 후)",
    "대안: ssh -L 8501:localhost:8501 wlsdud5035@34.64.124.90",
    "시나리오:",
    "1) 25구 위험도 지도 → Okabe-Ito 색맹 안전 팔레트 + 아이콘 (✅🔔⚠️🚨)",
    "2) 강남구 클릭 → 3-Layer 시계열",
    "3) 교차검증 탭 → 실측 JSON 동적 렌더 (0.643 등)",
    "4) 상관관계 탭 → 실측 Granger p-value 표시",
    "5) AI 리포트 탭 → 'AI 생성 · 인간 검토' 배지 + PDF 다운로드",
    "Plan B: 네트워크 죽으면 docs/images/screenshots/ PDF 시연",
])

# Slide 16 — 팀 · 남은 이슈 · Q&A
content(prs, "👥 팀 · 남은 이슈 · Q&A", [
    "5명 팀 + Claude Code 에이전트 시스템 (9개 모듈 CLAUDE.md)",
    "cd <모듈> && claude → 역할 특화 에이전트 · Stop hook 자동 메모리",
    "GitHub: CI 11 Job 🟢 · CODEOWNERS · 30 라벨 · 팀원 첫 이슈 #3-6",
    "남은 이슈: KDCA 실데이터(W18) · Kafka Consumer #4 · 실 DB 쿼리 #3",
    "         · Streamlit 스크린샷 #5 · Branch Protection #6",
    "리스크: 네이버 API 재판매 라이선스 · GPU 부재 · 팀원 역량 편차",
    "질문 환영합니다.",
])

OUT.parent.mkdir(parents=True, exist_ok=True); prs.save(str(OUT))
print(f"✅ {OUT.relative_to(OUT.parents[2])} ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} 슬라이드)")
