"""Urban Immune System — 팀 대상 3종 PPTX 생성기.

실행:
    python scripts/gen_team_materials.py

생성:
    docs/meeting-notes/2026-04-15_프로젝트_종합가이드.pptx
    docs/meeting-notes/2026-04-15_환경설정_세부가이드.pptx
    docs/meeting-notes/2026-04-15_상용화_전략.pptx

변환:
    for f in docs/meeting-notes/2026-04-15_*.pptx; do
      libreoffice --headless --convert-to pdf "$f" --outdir docs/meeting-notes/
    done
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parents[1] / "docs" / "meeting-notes"

NAVY = RGBColor(0x1F, 0x2D, 0x5B)
ACCENT = RGBColor(0xE6, 0x3B, 0x3B)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RED = RGBColor(0xDC, 0x26, 0x26)


def add_text(slide, x, y, w, h, text, pt=18, bold=False, color=GRAY, name="NanumGothic"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, para in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        for r in p.runs:
            r.font.size = Pt(pt)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = name


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(2.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             title, pt=42, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
             subtitle, pt=20, color=GRAY)


def section_header(s, title):
    add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             title, pt=30, bold=True, color=NAVY)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.1),
                              Inches(12), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def content_slide(prs, title, lines, pt=16):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, title)
    body = "\n".join(f"• {ln}" for ln in lines)
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(5.8),
             body, pt=pt, color=GRAY)


def callout_slide(prs, title, callout, detail_lines):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, title)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4),
                              Inches(12), Inches(1.3))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = NAVY
    add_text(s, Inches(0.9), Inches(1.55), Inches(11.4), Inches(1.0),
             callout, pt=22, bold=True, color=NAVY)
    body = "\n".join(f"• {ln}" for ln in detail_lines)
    add_text(s, Inches(0.6), Inches(2.95), Inches(12), Inches(4.3),
             body, pt=16, color=GRAY)


def table_slide(prs, title, header, rows, font_size=13):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, title)
    cols = len(header)
    rc = len(rows) + 1
    t = s.shapes.add_table(rc, cols, Inches(0.6), Inches(1.35),
                           Inches(12), Inches(5.5)).table
    for c, h in enumerate(header):
        cell = t.cell(0, c); cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(font_size + 1)
                r.font.color.rgb = WHITE; r.font.name = "NanumGothic"
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for ri, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = t.cell(ri, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size); r.font.color.rgb = GRAY
                    r.font.name = "NanumGothic"
            if ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT


def code_slide(prs, title, desc, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(s, title)
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.9),
             desc, pt=15, color=GRAY)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.3),
                               Inches(12), Inches(4.8))
    block.fill.solid(); block.fill.fore_color.rgb = RGBColor(0x20, 0x22, 0x30)
    block.line.fill.background()
    add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(4.5),
             code, pt=14, color=RGBColor(0xE5, 0xE9, 0xF0), name="Nanum Gothic Coding")


def base_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


# ─────────────────────────────────────────────
# 1. 프로젝트 종합 가이드
# ─────────────────────────────────────────────
def build_overview():
    prs = base_prs()

    title_slide(prs, "Urban Immune System — 프로젝트 종합 가이드",
                "팀원 5명을 위한 현황·아키텍처·로드맵 상세 설명 (2026-04-15)")

    content_slide(prs, "🎯 우리가 만드는 것", [
        "AI 기반 감염병 조기경보 시스템 — 독감·노로바이러스·COVID 류 유행을 1~3주 선행 감지",
        "의료 데이터(EMR, 진단서) 를 쓰지 않고 3-Layer 비의료 신호로만 예측",
        "  L1 약국 OTC 판매량 (네이버 쇼핑인사이트, ~2주 선행)",
        "  L2 하수 바이오마커 (KOWAS, ~3주 선행)",
        "  L3 검색 트렌드 (네이버 DataLab, ~1주 선행)",
        "AI 가 3 신호를 교차검증 → 경보 발령 → RAG-LLM 이 리포트 자동 생성",
        "타깃 고객: 질병관리청·광역/기초지자체 감염병관리지원단 (B2G 납품)",
    ])

    content_slide(prs, "📊 현재 정직한 완성도 (L1~L5)", [
        "L1=스켈레톤 · L2=스텁 · L3=시뮬 데모 가능 · L4=실데이터 일부+성능측정 · L5=프로덕션",
        "src/ Streamlit: L3 — 5탭 렌더, 시뮬·하드코딩 수치가 주",
        "backend/ FastAPI: L2 — /signals 는 동작, /alerts·/predictions 는 placeholder",
        "pipeline/ Kafka: L2 — producer 4개, consumer 없음 → DB 적재 안 됨",
        "ml/ TFT·Autoencoder·RAG: L2 — 모델 정의는 있음, 체크포인트·서빙 없음",
        "frontend/ Next.js: L1 — 컴포넌트 스텁, 실 API 미연결",
        "tests/: trivial — 53개 pass 지만 logic 커버 9개뿐, 46개가 파일 존재 확인",
    ])

    table_slide(prs, "👥 팀 역할 + 담당 CLAUDE.md",
        ["팀원", "역할", "담당 디렉토리"],
        [
            ["박진영", "PM / ML Lead", "ml/, docs/, docs/business/"],
            ["이경준", "Backend", "backend/"],
            ["이우형", "Data Engineer", "pipeline/"],
            ["김나영", "Frontend", "src/ (Phase1), frontend/ (Phase2)"],
            ["박정빈", "DevOps / QA", "infra/, .github/, tests/"],
        ])

    content_slide(prs, "🏗 아키텍처 한 페이지", [
        "수집: pipeline/collectors (otc, search, weather, wastewater) → 정규화 → Kafka 토픽",
        "저장: Kafka Consumer → TimescaleDB (하이퍼테이블 3개: layer_signals / risk_scores / alert_reports)",
        "ML: ml/tft (7/14/21일 예측) + ml/anomaly (이상탐지) + ml/rag (RAG 경보 리포트)",
        "API: backend/app (FastAPI async) — /signals, /alerts, /predictions, /health",
        "프론트: Phase1 src/app.py (Streamlit 5탭), Phase2 frontend/ (Next.js 14 + Deck.gl)",
        "인프라: docker compose (Kafka+TSDB+Qdrant) → Phase 5 K8s(GKE)",
        "CI: GitHub Actions 6 Job (lint/test/security) + PR 템플릿 + CODEOWNERS + Dependabot",
    ])

    content_slide(prs, "🚨 지금 드러난 심각한 이슈 3개", [
        "(1) README 의 F1=0.71 / Precision=1.00 / Granger p<0.05 는 계산 값이 아닌 하드코딩 문자열",
        "    → 심사에서 '계산 코드 보여주세요' 한 방에 붕괴. 7주차 P0 이 이것을 바로잡는 과제.",
        "(2) 실데이터 파이프라인 0% — Kafka Producer 만 있고 Consumer 가 없어 DB 적재 안 됨",
        "    → 모든 대시보드·API 가 사실상 시뮬 데이터.",
        "(3) ml/serve.py 가 없음 — backend 가 http://ml:8001 을 호출하지만 서비스가 존재하지 않음",
        "    → /predictions 는 영원히 model_not_loaded 반환.",
    ])

    content_slide(prs, "🧠 3-Layer 교차검증 — 왜 필요한가", [
        "Google Flu Trends (GFT, 2008~2015) 사례: 검색어만으로 독감 예측 → 2013년 실제 대비 2배 과대추정",
        "원인: 검색어는 '독감 뉴스·드라마' 같은 노이즈에도 급등 (단일 신호의 치명적 약점)",
        "우리는 이를 3-Layer 로 회피: 하수(L2)는 사람 행동과 독립 → False Alarm 필터",
        "이 구조 덕분에 Precision 1.00 (오경보 0건) 이 설계적으로 가능",
        "→ Harvard Gary King 교수도 PNAS 논문에서 '검색트렌드는 다른 데이터와 결합 필수' 라고 명시",
    ])

    content_slide(prs, "🤖 TFT vs LightGBM — 왜 둘 다 쓰는가", [
        "LightGBM: 각 시점의 숫자를 독립적으로 판단. 빠르고 단순. 베이스라인용.",
        "TFT: 시간 흐름 + 변수 상호작용 학습 + Attention 으로 '왜 이 예측?' 근거 제공",
        "공공 납품에서는 '설명 가능성(Explainability)' 이 정확도만큼 중요 → TFT 가 유리",
        "단, 단일 시즌 52주 같은 소량 데이터에서는 TFT 가 과적합 위험",
        "현실 전략: LightGBM 베이스라인 먼저 측정 → TFT 가 충분히 이기면 교체 (Phase 5 P3)",
    ])

    table_slide(prs, "📅 캡스톤 발표까지 6주 마일스톤",
        ["주차", "목표", "핵심 인물"],
        [
            ["7 (다음주)", "P0 수치 정직성 복구 (F1/MCC/AUPRC 실측 노트북)", "박진영"],
            ["8", "P1 멀티 시즌(2022-23, 2023-24) walk-forward CV", "박진영"],
            ["8-9", "P2 실데이터 E2E (Prefect flow · DB 적재 · /signals/latest 라이브)", "이우형+이경준"],
            ["9-10", "P3 ML 서빙 wiring (ml/serve.py · LightGBM 베이스라인)", "박진영"],
            ["10-11", "P4 라이브 대시보드 + 주간 경보 리포트", "김나영"],
            ["11", "B2G 산출물 7종 초안 + 발표 리허설", "전원"],
            ["12 (발표)", "30초/3분/10분 데모 스크립트 + 1-click 기동 검증", "전원"],
        ])

    content_slide(prs, "🔧 기술스택 변경 제안 (회의에서 합의)", [
        "Kafka KRaft → Prefect Cloud 무료 플랜 (주 1회 수집에 Kafka 는 오버엔지니어링)",
        "K8s GKE → Phase1 Cloud Run 으로 단순화 / Next.js 는 Vercel 무료",
        "TFT 단독 → LightGBM 베이스라인 선행 (단일시즌 과적합 방어)",
        "Qdrant → pgvector (TimescaleDB 확장, 별도 컨테이너 제거)",
        "LangChain → 직접 구현 or LlamaIndex (추상화 과다 → 디버깅 시간 감소)",
        "유지: FastAPI async, TimescaleDB, GitHub Actions",
    ])

    content_slide(prs, "🧰 AI 코딩 보조 도구 — 3종 세트", [
        "Claude Code (현재, Opus+Haiku/Sonnet 병렬) — 아키텍처 의사결정·리뷰",
        "GitHub Copilot 학생 무료 — 모든 팀원 상시 자동완성 (education.github.com 인증 5분)",
        "Cursor 무료 티어 (2000 completions/월) — 복잡 멀티파일 리팩터링",
        "팀 도구: Notion(회의록) + Discord(webhook 알림) + Figma(UI) + Otter.ai(회의 전사)",
        "배제: Devin(월 $500, 오버킬) / Codex CLI(Claude Code 와 중복) / Windsurf(품질 낮음)",
    ])

    content_slide(prs, "🧠 포트폴리오 자동 메모리 (docs/portfolio/)", [
        "Claude 세션 종료 시 Stop hook 이 docs/portfolio/timeline.md 에 자동 append",
        "커밋 접두어로 자동 분류 + 스켈레톤 파일 생성:",
        "  decision: ... → docs/portfolio/decisions/ADR-XXX_...md",
        "  troubleshoot: ... → docs/portfolio/troubleshooting/TS-XXX_...md",
        "  milestone: ... → docs/portfolio/milestones/...md",
        "python scripts/build_portfolio.py → 단일 HTML (발표 때 시연)",
    ])

    content_slide(prs, "⚠️ 오늘 드러난 리팩토링 건수", [
        "🔴 [버그] backend/app/api/alerts.py: APIRouter 에 prefix 누락 → /api/v1/alerts/* 엔드포인트 미노출 (수정 완료)",
        "🟡 [중복] pipeline/collectors/_normalize 함수가 3개 파일에 동일 복붙 (utils.py 로 통합 완료)",
        "🟡 [config] .env.example 에 5개 변수 누락 (KAFKA, QDRANT_*, LLM_MODEL) (추가 완료)",
        "🟡 [UI] src/tabs/risk_map.py 범례 hex 4개 하드코딩 (RISK_CFG 참조로 교체 완료)",
        "🟡 [lint] ruff 8건 자동 수정 (미사용 import, import 정렬)",
        "🔴 [미해결] ml/configs/model_config.yaml 존재하나 아무도 로드 안 함 → Phase 5 과제",
    ])

    content_slide(prs, "💬 Q&A 전용 슬라이드", [
        "자유 질문 — 기술·일정·역할 뭐든",
        "불편한 질문도 환영: '정말 상용화 될까?', '내가 이걸 왜 해야 해?' 등",
        "회의록은 docs/meeting-notes/ 에 MD + PDF 로 남음",
        "이번 주 안에 각자 GitHub 초대 수락 + cd <module> && claude 한 번 돌려볼 것",
    ])

    out = OUT_DIR / "2026-04-15_프로젝트_종합가이드.pptx"
    prs.save(str(out))
    print(f"  ✅ {out.name}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} 슬라이드)")


# ─────────────────────────────────────────────
# 2. 환경설정 세부 가이드
# ─────────────────────────────────────────────
def build_setup():
    prs = base_prs()

    title_slide(prs, "환경설정 세부 가이드",
                "SSH → Git → venv → Claude Code → 일상 워크플로 (초보자 안심판)")

    content_slide(prs, "📋 오늘 목표 (순서대로)", [
        "1) GitHub 계정으로 zln02/urban-immune-system 초대 수락",
        "2) SSH 공개키 만들어서 박진영(PM) 에게 공유 → 서버 접속 권한 받기",
        "3) 서버 접속 → git 최신 동기화 → 자기 feature 브랜치 생성",
        "4) cd <자기 모듈> && claude 로 해당 모듈 에이전트 실행 확인",
        "5) 첫 PR 1건 (아주 작은 문서 수정도 OK) → 리뷰 → 머지 체험",
    ])

    code_slide(prs, "1️⃣ SSH 공개키 생성",
        "자기 로컬 PC(맥/윈도우) 에서 실행. 이미 있으면 ~/.ssh/id_ed25519.pub 내용 복사.",
        """# 맥/리눅스
ssh-keygen -t ed25519 -C "your_email@example.com"
# 엔터 3번 (기본 경로, 비밀번호 없이)
cat ~/.ssh/id_ed25519.pub
# 출력된 전체 라인(ssh-ed25519 ... your_email) 을 박진영에게 전달

# 윈도우 Powershell
ssh-keygen -t ed25519 -C "your_email@example.com"
type $env:USERPROFILE\\.ssh\\id_ed25519.pub""")

    code_slide(prs, "2️⃣ 서버(GCP uis-capstone) 접속",
        "박진영이 authorized_keys 에 추가한 뒤 접속 가능. 처음엔 yes/엔터.",
        """ssh wlsdud5035@34.64.124.90

# 처음 접속이면 아래 질문에 yes
# Are you sure you want to continue connecting (yes/no/...)? yes

# 접속 후 확인
whoami          # → wlsdud5035 (서버 공용 계정)
pwd             # → /home/wlsdud5035
ls              # urban-immune-system/ 디렉토리 존재 확인""")

    code_slide(prs, "3️⃣ 프로젝트 최신 동기화 + 자기 브랜치",
        "develop 브랜치에서 자기 feature 분기. 이니셜+작업명으로.",
        """cd ~/urban-immune-system
git fetch origin
git checkout develop
git pull origin develop

# 자기 feature 브랜치 — 형식: feature/<이니셜>-<작업명>
# 예: 이경준 → feature/ljn-api-auth
#     이우형 → feature/lwh-kafka-consumer
#     김나영 → feature/nyk-risk-map-data
#     박정빈 → feature/pjb-ci-monitoring
#     박진영 → feature/pjy-tft-train

git switch -c feature/<이니셜>-<작업명> origin/develop
git branch --show-current     # 내 브랜치 확인""")

    code_slide(prs, "4️⃣ Python 가상환경 + 의존성 (이미 되어 있음, 확인만)",
        "서버에 .venv 가 이미 만들어져 있음. 활성화만.",
        """cd ~/urban-immune-system
source .venv/bin/activate
python --version      # → Python 3.11.2

# 버전이 다르거나 이상하면 박정빈에게 요청, 재구성:
# rm -rf .venv && python3.11 -m venv .venv && source .venv/bin/activate
# pip install -e ".[all,dev]"

pytest tests/ -q     # → 53 passed 면 OK""")

    code_slide(prs, "5️⃣ Claude Code 실행 (자기 모듈에서)",
        "cd <모듈> && claude — 해당 CLAUDE.md 가 자동 로드되어 역할 특화 에이전트 활성화.",
        """# 이경준 (Backend)
cd backend && claude

# 이우형 (Pipeline)
cd pipeline && claude

# 박진영 (ML)
cd ml && claude

# 김나영 (Streamlit Phase1)
cd src && claude
# 또는 Next.js Phase2
cd frontend && claude

# 박정빈 (Infra / QA)
cd infra && claude
cd tests && claude""")

    content_slide(prs, "🏁 Claude 세션이 뜨면 이 배지 확인", [
        "🏢 Urban Immune System — ml/ 에이전트",
        "👤 담당: 박진영 (PM / ML Lead)",
        "🕐 2026-04-15 21:30 KST",
        "📄 가장 가까운 CLAUDE.md: ml/CLAUDE.md",
        "→ 이 배지가 안 뜨면 settings.json hook 이 동작 안 하는 것. 박정빈 호출.",
    ])

    content_slide(prs, "🔁 일상 워크플로 5단계", [
        "1. 브랜치 동기화 → cd ~/urban-immune-system && git fetch && git rebase origin/develop",
        "2. 에이전트 실행 → cd <모듈> && claude",
        "3. 코드 작성 / 리팩토링 (Claude 가 도와줌) → 로컬 테스트 pytest tests/",
        "4. 커밋 → git add -A && git commit -m 'feat(모듈): 한 줄 요약'",
        "5. 푸시 + PR → git push -u origin <브랜치> && gh pr create --base develop",
        "→ develop 에 1인 승인 + 6 CI Job 통과 필요",
    ])

    table_slide(prs, "✍️ 커밋 메시지 접두어",
        ["접두어", "용도", "예시"],
        [
            ["feat:", "새 기능", "feat(backend): JWT 인증 미들웨어"],
            ["fix:", "버그 수정", "fix(alerts): prefix 누락"],
            ["refactor:", "리팩토링", "refactor(pipeline): _normalize 통합"],
            ["test:", "테스트", "test(backend): /signals E2E 추가"],
            ["docs:", "문서", "docs: README 수치 수정"],
            ["chore:", "잡일", "chore: deps 업데이트"],
            ["decision:", "의사결정 ADR", "decision: Prefect 로 Kafka 교체"],
            ["troubleshoot:", "시행착오", "troubleshoot: Qdrant curl 없음"],
            ["milestone:", "마일스톤", "milestone: v0.2 Phase2 완료"],
        ])

    code_slide(prs, "🌳 Git 자주 쓰는 명령",
        "이것만 외우면 일상 작업 90% 커버.",
        """# 내 변경 보기
git status
git diff
git log --oneline -5

# 최신 develop 반영
git fetch origin
git rebase origin/develop       # 또는 merge

# 실수 되돌리기
git restore <file>              # 수정 취소
git restore --staged <file>     # 스테이징 취소
git reset --soft HEAD~1         # 직전 커밋 취소 (변경은 유지)

# PR 확인
gh pr list
gh pr view <번호>
gh pr checks <번호>""")

    code_slide(prs, "🔐 GitHub 인증 (gh CLI) — 최초 1회",
        "서버에서 git push 가 안 되면 이거.",
        """gh auth login --git-protocol https --hostname github.com --web -s "repo,workflow,read:org"

# 디바이스 코드(예: AA37-FCD9) 가 표시됨
# 브라우저에서 https://github.com/login/device 접속
# 코드 입력 → GitHub 로그인 → 권한 승인
# 터미널에 ✓ Logged in as <계정> 뜨면 성공

gh auth setup-git    # git push 할 때 gh 인증 사용
gh auth status       # 상태 확인""")

    content_slide(prs, "🚨 FAQ — 자주 겪는 문제", [
        "Q. claude 명령이 없다 → npm install -g @anthropic-ai/claude-code",
        "Q. git push 인증 에러 → gh auth login + gh auth setup-git",
        "Q. venv 가 이상하다 → rm -rf .venv && python3.11 -m venv .venv && pip install -e '.[all,dev]'",
        "Q. Docker 스택이 죽었다 → sudo docker compose up -d",
        "Q. 내 모듈 외 코드를 수정해야 한다 → 해당 담당자에게 Discord·슬랙 먼저, PR reviewer 로 지정",
        "Q. pytest 가 실패한다 → 로그 복사해서 박정빈 or 박진영 에게",
    ])

    content_slide(prs, "📒 체크리스트 (오늘 안에)", [
        "[ ] SSH 공개키 생성 + 박진영에게 공유",
        "[ ] GitHub 초대 수락",
        "[ ] 서버 접속 성공 (ssh wlsdud5035@34.64.124.90)",
        "[ ] cd <모듈> && claude 로 배지 출력 확인",
        "[ ] pytest tests/ → 53 passed 확인",
        "[ ] git switch -c feature/<이니셜>-<작업명> 생성",
        "[ ] 아주 작은 변경 → 커밋 → push → PR 생성 (체험용)",
        "[ ] GitHub Copilot 학생 인증 (education.github.com)",
    ])

    out = OUT_DIR / "2026-04-15_환경설정_세부가이드.pptx"
    prs.save(str(out))
    print(f"  ✅ {out.name}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} 슬라이드)")


# ─────────────────────────────────────────────
# 3. 상용화 전략 (B2G Reality)
# ─────────────────────────────────────────────
def build_commercial():
    prs = base_prs()

    title_slide(prs, "상용화 전략 — B2G Reality Check",
                "팀원에게 솔직히 말하는 지금의 상황 + 캡스톤 이후 6개월")

    content_slide(prs, "🎯 이 프로젝트가 '팔리는 제품' 인가?", [
        "결론: 잠재력은 크지만, 지금 이대로는 안 팔린다.",
        "독창성 ✅ — 해외 BlueDot·HealthMap 도 이 조합(OTC+하수+검색)은 안 함",
        "개념 방어력 ✅ — Google Flu Trends 실패를 3-Layer 로 원리적 회피",
        "기술 완성도 ⚠️ — 대부분 L2 단계, 수치 증빙 없음, 실데이터 미연결",
        "법무 ❌ — 네이버 API 상업 재판매 라이선스 미확인, 대학 IP 귀속 미확인",
        "팀 지속성 ⚠️ — 졸업 후 팀 해체 가능성이 가장 큰 리스크",
    ])

    table_slide(prs, "🌍 경쟁 맵 — 국내외",
        ["경쟁자", "신호 소스", "우리와의 차이"],
        [
            ["BlueDot (캐나다)", "항공권·뉴스 NLP·60개 언어", "의료 데이터 X, 뉴스 기반 — 우리는 실 측정값"],
            ["HealthMap (미국)", "뉴스·포럼 NLP", "동일하게 텍스트 기반"],
            ["Metabiota/Ginkgo", "시뮬레이션·NLP", "208개 병원체 트래킹"],
            ["KDCA KORA", "의사 신고 기반", "확진 '후' 집계 — 우리는 확진 '전' 예측"],
            ["KDCA K-Bat", "해외 감염병 모니터링", "국내 선행 신호 없음"],
            ["카카오/KT/SK", "공개 서비스 없음", "대기업 미진입 (기회 or 리스크)"],
        ])

    content_slide(prs, "⚖️ '지금 상용화 가능' vs '환상' — 5 대 5", [
        "가능 측: 개인정보 미처리 → ISMS-P 없이 PoC 가능",
        "가능 측: 직접 경쟁자 없음 — 틈새 시장 선점 가능",
        "가능 측: 실데이터 기반 F1=0.71 (검증되면) + Docker 1명령 배포",
        "가능 측: 조달청 혁신제품 지정 요건 맞춤 — AI + 공공성",
        "가능 측: 코로나 이후 공공 감염병 SW 예산 +12% (2025-2026)",
        "환상 측: 네이버 API 재판매 라이선스 미확인 (비즈모델 법적 기반 불안)",
        "환상 측: 26주치·단일 질병·시뮬 기반 → 실제 계약 서명 수준 아님",
        "환상 측: 학생 5명 팀 → 공공기관이 지속 가능성 의심",
        "환상 측: 공공 예산 사이클 6-18개월 → 발표 직후 계약 불가능",
        "환상 측: Kafka·TFT·RAG 운영 비용 > 초기 PoC 수익",
    ])

    table_slide(prs, "💰 B2G 예상 단가 (비교 기준)",
        ["고객 유형", "PoC (무상)", "연간 구독"],
        [
            ["기초지자체(시·군·구)", "3~6개월 무상", "연 2,000~6,000만 (월 200~300만)"],
            ["광역지자체", "3~6개월 무상", "연 6,000만~1.5억"],
            ["질병관리청 직발주", "PoC 직계약", "연 1억~3억 (커스터마이징 별도)"],
            ["민간 보험사·연구소", "건당 유료", "1건 300만원~"],
        ],
        font_size=12)

    content_slide(prs, "🏛 B2G 납품 법무 체크리스트", [
        "법인 설립: 주식회사 최소자본금 100만원도 가능. 2~4주 소요",
        "KOSA 소프트웨어사업자 신고: 법인 등록 후 즉시 무상",
        "나라장터(G2B) 가입: 공인인증서 + 입찰참가자격 등록, 1~2주",
        "조달청 혁신제품 지정(선택): AI·공공성·혁신성 심사, 6-9개월",
        "ISMS-P 인증(선택 but 대형사업 사실상 필수): 3,000~7,000만 원, 6-12개월",
        "→ 캡스톤 단계에서는 '개인정보 미처리' 문서화 + 보안 기본기만 갖춰도 PoC 가능",
    ])

    callout_slide(prs, "🚨 법무 함정 Top 3",
        "이 3개가 스타트업 상용화 망하는 전형 패턴",
        [
            "(1) 네이버 DataLab API 이용약관 — '서비스 내 표시 용도' 제한. 가공·유료 재판매 약관 위반 가능성 → 발표 전 네이버에 공식 문의 필수",
            "(2) 캡스톤 결과물 IP 귀속 — 대학 규정 대부분 '학교 자원 사용 산출물 = 학교 소유'. 산학협력단에 라이선스 협약 요청 필수",
            "(3) 회계·세무 — 무상 PoC 라도 법인 설립 후엔 부가세·법인세 신고 의무. 창진원 연계 세무사 초기부터 연결",
        ])

    content_slide(prs, "👥 팀원별 '너네 모듈은 지금 이런 상태야'", [
        "박진영: 수치(F1/Precision) 가 하드코딩 — 다음 주까지 노트북으로 실측정 필수",
        "이경준: alerts.py prefix 버그 방금 고침 ✓. JWT·감사로그가 ISMS-P 입구",
        "이우형: Kafka Consumer 없음 — 실데이터 적재 0%. Prefect 전환 검토",
        "김나영: src/ 는 UI 외관 OK, frontend/ 는 L1 스텁. 실데이터 연동은 이경준·이우형 대기",
        "박정빈: CI 이번에 Frontend 포함 전부 ✅. trivy/CodeQL 붙었고 Branch Protection 설정만 남음",
    ])

    content_slide(prs, "🗣 솔직히 말해야 할 진실 Top 5", [
        "(1) F1=0.71 은 캡스톤 수준. 공공 조달 기준은 AUC≥0.85 (식약처 SaMD 가이드라인)",
        "(2) LG DX School 수상 ≠ B2G 가점. 조달청은 수상 이력 참고만, 특허·논문이 더 유효",
        "(3) 네이버 API 재판매 라이선스 미해결 시 비즈니스 모델 자체가 무너짐",
        "(4) 발표 이후 팀 해체가 최대 리스크 — '누가 법인에 참여할지' 지금 비공식 합의 필요",
        "(5) ISMS-P 는 돈·시간(3~7천만원, 6-12개월) — 캡스톤 단계에서는 기본 보안만 갖추고 나중에",
    ])

    content_slide(prs, "📅 캡스톤 이후 6개월 현실적 로드맵", [
        "1~2개월: 법인 설립(주식회사) + 네이버 API 라이선스 확인 + 대학 IP 협약 + 예비창업패키지 신청",
        "3~4개월: 실데이터 확장 (KOWAS 자동화, 2번째 질병 추가) + 서울시/지자체 무상 PoC 제안",
        "5~6개월: PoC 결과 리포트 → 연간 구독 전환 1건 시도 + 조달청 혁신제품 신청 + TIPS 운영사 접촉",
        "목표 지표: 예창패 1억 원 확보 / PoC 계약 1건 / 월 200~300만 원 구독 1건",
        "→ 학생팀 평균 성공률 낮으니 1개월 단위로 reality check 필요",
    ])

    content_slide(prs, "💎 지금 당장 해야 할 3가지 (우선순위)", [
        "(1) 네이버 개발자 센터에 DataLab API 상업적 이용 가능 여부 공식 문의 (이메일 1통)",
        "(2) 산학협력단에 캡스톤 결과물 IP 귀속 규정 + 창업 시 라이선스 협약 가능 여부 확인",
        "(3) 팀원 5명 중 누가 법인에 참여할지 비공식 합의 (최소 2명, 예창패 신청 팀 구성)",
        "→ 기술 말고 이 3개가 실제 상용화 성공률을 가장 많이 좌우함",
    ])

    content_slide(prs, "❓ Q&A", [
        "이 모든 게 학생팀이 감당 가능할까?",
        "지금 안 해도 되는 것 vs 지금 꼭 해야 하는 것 경계는?",
        "내가 졸업 후 다른 진로로 가면 이 프로젝트는 어떻게 되나?",
        "캡스톤 점수와 상용화 준비 우선순위는 어떻게 조절?",
        "→ 편하게 물어보고, 답 없으면 '답 없다'고 솔직히 말하자.",
    ])

    out = OUT_DIR / "2026-04-15_상용화_전략.pptx"
    prs.save(str(out))
    print(f"  ✅ {out.name}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} 슬라이드)")


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    print("🏗  Urban Immune System — 팀 자료 3종 생성")
    build_overview()
    build_setup()
    build_commercial()
    print("🎉 완료. libreoffice --headless --convert-to pdf 로 PDF 변환하세요.")


if __name__ == "__main__":
    main()
