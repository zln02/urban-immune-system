"""Slide 3·16 오버레이 패치본 생성 (이미지 기반 PPTX용).

원본 PPTX는 모든 슬라이드가 단일 PICTURE 셰이프(Gamma/Tome export 추정).
텍스트 직접 수정 불가하므로 이미지 위에 작은 회색 텍스트박스를 얹는다.

출력: docs/slides/2026-04-30_중간발표_v2.pptx (원본은 그대로 유지)
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

ROOT = Path(__file__).parent.parent
SRC = ROOT / "docs" / "slides" / "2026-04-30_중간발표.pptx"
DST = ROOT / "docs" / "slides" / "2026-04-30_중간발표_v2.pptx"

p = Presentation(str(SRC))
SW, SH = p.slide_width, p.slide_height  # 12192000 x 6858000


def add_label(slide, text: str, *, top_emu: int, font_pt: float = 11,
              color=RGBColor(0x64, 0x74, 0x8B)) -> None:
    """슬라이드 하단 가운데에 회색 캡션 텍스트박스를 추가."""
    box_w = Emu(int(SW * 0.85))
    box_h = Emu(Pt(font_pt * 1.6).emu)
    left = Emu(int((SW - int(box_w)) / 2))
    tx = slide.shapes.add_textbox(left, Emu(top_emu), box_w, box_h)
    tf = tx.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = 2  # center
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.color.rgb = color
    run.font.italic = True
    # 배경색 살짝 반투명한 흰색 박스
    fill = tx.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    line = tx.line
    line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    line.width = Emu(6350)  # 0.5pt


# Slide 3 (idx=2): 합성가정 라벨
add_label(
    p.slides[2],
    "※ 위 선행시간은 합성데이터 모델링 가정값. 실측은 Slide 12 참조 — L1=8주 / L2=2주 / L3=3주 / 17지역 평균 5.9주.",
    top_emu=int(SH * 0.92),
    font_pt=10,
)

# Slide 16 (idx=15): KCDC 출처 라벨
add_label(
    p.slides[15],
    "출처: KCDC 감염병포털 (infpublic.kdca.go.kr) 2025-W49 confirmed_cases · pipeline/collectors/kcdc_collector.py",
    top_emu=int(SH * 0.93),
    font_pt=9,
)

p.save(str(DST))
print(f"OK → {DST}  ({DST.stat().st_size / 1024:.1f}KB)")
print("원본 보존: docs/slides/2026-04-30_중간발표.pptx")
